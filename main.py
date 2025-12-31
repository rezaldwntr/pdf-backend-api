# -*- coding: utf-8 -*-
"""
Aplikasi Web Konverter PDF (c) 2024
Versi: 6.0 (Merge, Split, & Smart Compression with Target Size)
"""
import os
import shutil
import logging
import tempfile
import io
from typing import List, Optional
from enum import Enum
from zipfile import ZipFile

# Framework
from fastapi import FastAPI, File, UploadFile, HTTPException, BackgroundTasks, Form
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware

# Library Konversi
from pdf2docx import Converter
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import pdfplumber
import pandas as pd
from openpyxl.styles import Border, Side, Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter

# Konfigurasi Ukuran File (25MB)
MAX_FILE_SIZE = 25 * 1024 * 1024 

app = FastAPI(
    title="Aplikasi Konverter PDF",
    description="API untuk mengubah format file dari PDF ke format lainnya.",
    version="6.0",
)

# === KONFIGURASI CORS ===
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- HELPER FUNCTIONS ---
def cleanup_folder(path: str):
    try:
        if os.path.exists(path):
            shutil.rmtree(path)
            logging.info(f"Deleted temp folder: {path}")
    except Exception as e:
        logging.error(f"Error cleaning up: {e}")

def validate_file(file: UploadFile):
    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="File harus format PDF")
    file.file.seek(0, 2)
    file_size = file.file.tell()
    file.file.seek(0)
    if file_size > MAX_FILE_SIZE:
        raise HTTPException(status_code=400, detail=f"Ukuran file terlalu besar (Maks {MAX_FILE_SIZE/1024/1024}MB)")

@app.get("/")
def read_root():
    return {"message": "Server PDF Backend (V6.0 Complete Suite) is Running!"}

# ==========================================
# 1. KONVERSI DOKUMEN (Word, Excel, PPT, Image)
# ==========================================

# ... [FITUR PDF TO DOCX] ...
@app.post("/convert/pdf-to-docx")
def convert_pdf_to_docx(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    validate_file(file)
    tmp_dir = tempfile.mkdtemp()
    tmp_pdf_path = os.path.join(tmp_dir, file.filename)
    docx_filename = os.path.splitext(file.filename)[0] + ".docx"
    tmp_docx_path = os.path.join(tmp_dir, docx_filename)
    try:
        with open(tmp_pdf_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        cv = Converter(tmp_pdf_path)
        cv.convert(tmp_docx_path, start=0, end=None, multiprocess=False)
        cv.close()
        background_tasks.add_task(cleanup_folder, tmp_dir)
        return FileResponse(path=tmp_docx_path, filename=docx_filename, media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
    except Exception as e:
        cleanup_folder(tmp_dir)
        raise HTTPException(status_code=500, detail=f"Gagal convert Word: {str(e)}")

# ... [FITUR PDF TO EXCEL] ...
@app.post("/convert/pdf-to-excel")
def convert_pdf_to_excel(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    validate_file(file)
    tmp_dir = tempfile.mkdtemp()
    tmp_pdf_path = os.path.join(tmp_dir, file.filename)
    xlsx_filename = os.path.splitext(file.filename)[0] + ".xlsx"
    tmp_xlsx_path = os.path.join(tmp_dir, xlsx_filename)
    try:
        with open(tmp_pdf_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        with pd.ExcelWriter(tmp_xlsx_path, engine='openpyxl') as writer:
            with pdfplumber.open(tmp_pdf_path) as pdf:
                writer.book.create_sheet("Hasil Konversi")
                worksheet = writer.book["Hasil Konversi"]
                if "Sheet" in writer.book.sheetnames: del writer.book["Sheet"]
                current_row = 1
                thin_border = Border(left=Side(style='thin'), right=Side(style='thin'), top=Side(style='thin'), bottom=Side(style='thin'))
                
                # Helper internal untuk excel
                def is_inside(bbox, tables):
                    mx, my = (bbox[0]+bbox[2])/2, (bbox[1]+bbox[3])/2
                    for t in tables:
                        if t[0]<=mx<=t[2] and t[1]<=my<=t[3]: return True
                    return False
                
                for page in pdf.pages:
                    tables = page.find_tables()
                    table_bboxes = [t.bbox for t in tables]
                    words = page.extract_words()
                    non_table_words = [w for w in words if not is_inside((w['x0'],w['top'],w['x1'],w['bottom']), table_bboxes)]
                    
                    # Simple text clustering (simplified for brevity in this merged version)
                    non_table_text = []
                    if non_table_words:
                         non_table_text = sorted(non_table_words, key=lambda x: x['top'])
                    
                    # Write simple text
                    for w in non_table_text:
                        worksheet.cell(row=current_row, column=1, value=w['text'])
                        current_row += 1
                    
                    # Write tables
                    for t in tables:
                        data = t.extract()
                        if data:
                            df = pd.DataFrame(data)
                            for r in dataframe_to_rows(df, index=False, header=False):
                                for c_idx, val in enumerate(r, 1):
                                    c = worksheet.cell(row=current_row, column=c_idx, value=val)
                                    c.border = thin_border
                                current_row += 1
                        current_row += 1
                    current_row += 2
        
        background_tasks.add_task(cleanup_folder, tmp_dir)
        return FileResponse(path=tmp_xlsx_path, filename=xlsx_filename, media_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
    except Exception as e:
        cleanup_folder(tmp_dir)
        # Fallback simple jika logic kompleks error (agar tidak 500 terus)
        logging.error(f"Excel error: {e}")
        raise HTTPException(status_code=500, detail="Gagal convert Excel. Pastikan file tidak corrupt.")

# Helper function untuk excel rows (pengganti utils pandas)
def dataframe_to_rows(df, index=False, header=False):
    if header: yield df.columns.tolist()
    for row in df.itertuples(index=index, name=None): yield row

# ... [FITUR PDF TO PPT - V5.0 Ultimate] ...
@app.post("/convert/pdf-to-ppt")
def convert_pdf_to_ppt(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    validate_file(file)
    tmp_dir = tempfile.mkdtemp()
    tmp_pdf_path = os.path.join(tmp_dir, file.filename)
    ppt_filename = os.path.splitext(file.filename)[0] + ".pptx"
    tmp_ppt_path = os.path.join(tmp_dir, ppt_filename)
    try:
        with open(tmp_pdf_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        prs = Presentation()
        doc = fitz.open(tmp_pdf_path)
        if len(doc) > 0:
            p1 = doc[0]
            prs.slide_width = int((p1.rect.width / 72) * 914400)
            prs.slide_height = int((p1.rect.height / 72) * 914400)
        
        for page in doc:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            # Images
            img_blocks = [b for b in page.get_text("dict", flags=fitz.TEXT_PRESERVE_IMAGES)["blocks"] if b['type']==1]
            for b in img_blocks:
                img_path = os.path.join(tmp_dir, f"img_{os.urandom(4).hex()}.{b['ext']}")
                with open(img_path, "wb") as f: f.write(b["image"])
                x0,y0,x1,y1 = b["bbox"]
                try: slide.shapes.add_picture(img_path, Inches(x0/72), Inches(y0/72), width=Inches((x1-x0)/72), height=Inches((y1-y0)/72))
                except: pass
            
            # Text
            text_blocks = [b for b in page.get_text("dict")["blocks"] if b['type']==0]
            for b in text_blocks:
                for line in b["lines"]:
                    lx0,ly0,lx1,ly1 = line["bbox"]
                    txBox = slide.shapes.add_textbox(Inches(lx0/72), Inches(ly0/72), Inches((lx1-lx0)/72), Inches((ly1-ly0)/72))
                    tf = txBox.text_frame
                    tf.word_wrap = False
                    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                    p = tf.paragraphs[0]
                    for span in line["spans"]:
                        if not span["text"].strip(): continue
                        run = p.add_run()
                        run.text = span["text"]
                        run.font.size = Pt(span["size"])
                        try:
                            c = span["color"]
                            run.font.color.rgb = RGBColor((c>>16)&0xFF, (c>>8)&0xFF, c&0xFF)
                        except: pass
                        if span["flags"] & 16: run.font.bold = True
                        if span["flags"] & 2: run.font.italic = True
        doc.close()
        prs.save(tmp_ppt_path)
        background_tasks.add_task(cleanup_folder, tmp_dir)
        return FileResponse(path=tmp_ppt_path, filename=ppt_filename, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    except Exception as e:
        cleanup_folder(tmp_dir)
        raise HTTPException(status_code=500, detail=f"Gagal convert PPT: {str(e)}")

# ... [FITUR PDF TO IMAGE] ...
@app.post("/convert/pdf-to-image")
def convert_pdf_to_image(background_tasks: BackgroundTasks, file: UploadFile = File(...), output_format: str = "png"):
    validate_file(file)
    tmp_dir = tempfile.mkdtemp()
    tmp_pdf_path = os.path.join(tmp_dir, file.filename)
    zip_filename = os.path.splitext(file.filename)[0] + "_images.zip"
    tmp_zip_path = os.path.join(tmp_dir, zip_filename)
    try:
        with open(tmp_pdf_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        doc = fitz.open(tmp_pdf_path)
        with ZipFile(tmp_zip_path, 'w') as zipf:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=200)
                img_name = f"page_{i+1}.{output_format}"
                img_path = os.path.join(tmp_dir, img_name)
                pix.save(img_path, output="jpg" if output_format.lower() in ['jpg','jpeg'] else output_format)
                zipf.write(img_path, img_name)
        doc.close()
        background_tasks.add_task(cleanup_folder, tmp_dir)
        return FileResponse(path=tmp_zip_path, filename=zip_filename, media_type='application/zip')
    except Exception as e:
        cleanup_folder(tmp_dir)
        raise HTTPException(status_code=500, detail=f"Gagal convert Image: {str(e)}")


# ==========================================
# 2. FITUR ALAT (MERGE, SPLIT, COMPRESS)
# ==========================================

# === FITUR 5: GABUNGKAN PDF (MERGE) ===
@app.post("/tools/merge-pdf")
def merge_pdf(background_tasks: BackgroundTasks, files: List[UploadFile] = File(...)):
    if len(files) < 2:
        raise HTTPException(status_code=400, detail="Minimal upload 2 file PDF.")
    
    tmp_dir = tempfile.mkdtemp()
    merged_filename = "merged_document.pdf"
    tmp_merged_path = os.path.join(tmp_dir, merged_filename)

    try:
        merged_doc = fitz.open()
        for file in files:
            if not file.filename.lower().endswith(".pdf"): continue
            file_path = os.path.join(tmp_dir, file.filename)
            with open(file_path, "wb") as f:
                shutil.copyfileobj(file.file, f)
            doc = fitz.open(file_path)
            merged_doc.insert_pdf(doc)
            doc.close()
        
        merged_doc.save(tmp_merged_path)
        merged_doc.close()
        background_tasks.add_task(cleanup_folder, tmp_dir)
        return FileResponse(path=tmp_merged_path, filename=merged_filename, media_type='application/pdf')
    except Exception as e:
        cleanup_folder(tmp_dir)
        raise HTTPException(status_code=500, detail=f"Gagal Merge: {str(e)}")

# === FITUR 6: PISAHKAN PDF (SPLIT) ===
@app.post("/tools/split-pdf")
def split_pdf(background_tasks: BackgroundTasks, pages: str = Form(...), file: UploadFile = File(...)):
    validate_file(file)
    tmp_dir = tempfile.mkdtemp()
    tmp_pdf_path = os.path.join(tmp_dir, file.filename)
    split_filename = f"split_{file.filename}"
    tmp_split_path = os.path.join(tmp_dir, split_filename)

    try:
        with open(tmp_pdf_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        src_doc = fitz.open(tmp_pdf_path)
        new_doc = fitz.open()
        
        # Parse range "1,3,5-7"
        selected_indices = []
        try:
            for part in pages.split(','):
                if '-' in part:
                    s, e = map(int, part.split('-'))
                    selected_indices.extend(range(s-1, e))
                else:
                    selected_indices.append(int(part)-1)
        except:
            raise HTTPException(status_code=400, detail="Format halaman salah. Contoh: 1,3,5-10")
        
        for idx in selected_indices:
            if 0 <= idx < len(src_doc):
                new_doc.insert_pdf(src_doc, from_page=idx, to_page=idx)
        
        if len(new_doc) == 0:
            raise HTTPException(status_code=400, detail="Halaman tidak ditemukan/kosong.")

        new_doc.save(tmp_split_path)
        src_doc.close(); new_doc.close()
        background_tasks.add_task(cleanup_folder, tmp_dir)
        return FileResponse(path=tmp_split_path, filename=split_filename, media_type='application/pdf')
    except Exception as e:
        cleanup_folder(tmp_dir)
        raise HTTPException(status_code=500, detail=f"Gagal Split: {str(e)}")

# === FITUR 7: KOMPRES PDF (ADVANCED) ===
class CompressionType(str, Enum):
    RECOMMENDED = "recommended"
    TARGET = "target"

@app.post("/tools/compress-pdf")
def compress_pdf(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    compression_type: CompressionType = Form(CompressionType.RECOMMENDED),
    target_size_kb: Optional[int] = Form(None)
):
    validate_file(file)
    tmp_dir = tempfile.mkdtemp()
    tmp_pdf_path = os.path.join(tmp_dir, file.filename)
    comp_filename = f"compressed_{file.filename}"
    tmp_comp_path = os.path.join(tmp_dir, comp_filename)

    try:
        with open(tmp_pdf_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        # Ukuran Awal
        original_size = os.path.getsize(tmp_pdf_path)
        logging.info(f"Original size: {original_size} bytes")

        doc = fitz.open(tmp_pdf_path)

        # LOGIKA 1: KOMPRESI REKOMENDASI (Standard Deflate)
        if compression_type == CompressionType.RECOMMENDED:
            doc.save(
                tmp_comp_path,
                garbage=4,  # Hapus unused objects
                deflate=True, # Kompres stream
                clean=True
            )
        
        # LOGIKA 2: KOMPRESI TARGET SIZE (Best Effort)
        elif compression_type == CompressionType.TARGET and target_size_kb:
            target_bytes = target_size_kb * 1024
            
            # Tahap 1: Coba kompresi standar dulu
            doc.save(tmp_comp_path, garbage=4, deflate=True)
            current_size = os.path.getsize(tmp_comp_path)
            
            # Jika masih terlalu besar, lakukan Downsampling Gambar Agresif
            if current_size > target_bytes:
                logging.info("Standard compression not enough, starting aggressive downsampling...")
                
                # Kita coba 2 level: 96 DPI dan 72 DPI
                for dpi_level in [96, 72, 50]: # Loop penurunan kualitas
                    if current_size <= target_bytes:
                        break # Sudah cukup
                    
                    # Buat PDF baru
                    new_doc = fitz.open()
                    
                    # Render setiap halaman jadi gambar (Rasterize) lalu masukkan ke PDF baru
                    # Ini cara paling ampuh mengecilkan ukuran file scan/foto
                    for page in doc:
                        pix = page.get_pixmap(dpi=dpi_level)
                        # Kompresi gambar JPG
                        img_bytes = pix.pil_tobytes(format="JPEG", quality=70, optimize=True)
                        
                        img_page = new_doc.new_page(width=page.rect.width, height=page.rect.height)
                        img_page.insert_image(page.rect, stream=img_bytes)
                    
                    new_doc.save(tmp_comp_path, garbage=4, deflate=True)
                    new_doc.close()
                    
                    current_size = os.path.getsize(tmp_comp_path)
                    logging.info(f"Tried DPI {dpi_level}, new size: {current_size}")

        else:
            # Fallback default
            doc.save(tmp_comp_path, garbage=4, deflate=True)

        doc.close()
        
        final_size = os.path.getsize(tmp_comp_path)
        logging.info(f"Final size: {final_size} bytes")

        background_tasks.add_task(cleanup_folder, tmp_dir)
        return FileResponse(path=tmp_comp_path, filename=comp_filename, media_type='application/pdf')

    except Exception as e:
        cleanup_folder(tmp_dir)
        logging.error(f"ERROR COMPRESS: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Gagal kompres PDF: {str(e)}")
